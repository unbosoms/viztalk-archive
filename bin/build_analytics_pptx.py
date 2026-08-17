#!/usr/bin/env python3
"""analytics/README.md の内容を PowerPoint (.pptx) にする。

出力: analytics/vizトーク-アーカイブ-分析用データセット.pptx

実行:
  .venv/bin/python3 bin/build_analytics_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = ROOT / "analytics" / "vizトーク-アーカイブ-分析用データセット.pptx"

# --- Vizトーク ブランドカラー ---
TEAL = RGBColor(0x3B, 0x8A, 0x99)
TEAL_DARK = RGBColor(0x2A, 0x63, 0x6E)
ORANGE = RGBColor(0xE8, 0x82, 0x3B)
CREAM = RGBColor(0xFA, 0xF3, 0xE7)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCC, 0xCC, 0xCC)

FONT_JP = "Hiragino Sans"  # macOS 標準
FONT_MONO = "Menlo"


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_JP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_shape(slide, shape, x, y, w, h, *, fill=None, line=None, line_width=1.0):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is not None:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_slide_header(slide, title, subtitle=None):
    """統一ヘッダー: 左肩に細いtealバー + タイトル + サブ"""
    # tealバー
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5),
              fill=TEAL)
    # ページ上部ライン
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0), Inches(12.98), Inches(0.08),
              fill=ORANGE)
    add_text(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=TEAL_DARK)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GRAY)


def add_footer(slide, text="Vizトーク アーカイブ 分析用データセット"):
    add_text(slide, Inches(0.7), Inches(7.15), Inches(12), Inches(0.3),
             text, size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# =========================================================
# スライド構築
# =========================================================

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, CREAM)
    # ロゴエリア (シンプルな飾り)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.15), Inches(1.8), Inches(1.0), Inches(1.0),
              fill=TEAL)
    add_text(slide, Inches(6.15), Inches(1.8), Inches(1.0), Inches(1.0),
             "V", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0), Inches(3.1), Inches(13.33), Inches(0.9),
             "Vizトーク アーカイブ", size=44, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(3.95), Inches(13.33), Inches(0.6),
             "分析用データセット", size=28, bold=False, color=ORANGE, align=PP_ALIGN.CENTER)

    # 装飾ライン
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.665), Inches(4.75), Inches(2.0), Inches(0.03),
              fill=ORANGE)

    add_text(slide, Inches(0), Inches(5.0), Inches(13.33), Inches(0.4),
             "第36回〜最新回 / 5 CSV / スタースキーマ設計", size=14, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0), Inches(6.9), Inches(13.33), Inches(0.3),
             "Vizトーク 主催者陣 配布用 (2026年8月)", size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_slide_header(slide, "データセット概要", "Vizトーク 全放送メタデータの分析用抽出")

    items = [
        ("配布範囲", "現在は主催者・スピーカー陣への個別配布のみ (git管理外)"),
        ("収録範囲", "第36回 (2023年8月) 〜 最新の収録済み回、計110回"),
        ("粒度", "回単位のメトリクス + チャプター単位の詳細"),
        ("生成元", "スペース音源 → Whisper文字起こし → LLM でチャプター/トピック抽出 → #Vizトーク 実況ツイート数"),
        ("フォーマット", "UTF-8 CSV × 5ファイル、スタースキーマ設計"),
        ("推奨ツール", "Tableau / Power BI / Looker Studio 等の BI ツール"),
    ]

    y = 1.7
    for label, desc in items:
        # 色帯
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(0.1), Inches(0.6),
                  fill=TEAL)
        add_text(slide, Inches(1.0), Inches(y), Inches(2.2), Inches(0.6),
                 label, size=15, bold=True, color=TEAL_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(3.3), Inches(y), Inches(9.7), Inches(0.6),
                 desc, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.75

    add_footer(slide)


def slide_files_list(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_slide_header(slide, "ファイル一覧", "全て ep + date を join key に連結")

    files = [
        ("episodes.csv", "110", "Fact テーブル", "1行=1回。全メトリクスがここに集約"),
        ("episode_speakers.csv", "286", "Dim (long)", "1行=1回×1スピーカー"),
        ("episode_topics.csv", "4,313", "Dim (long)", "1行=1回×1トピック"),
        ("chapters.csv", "2,562", "Dim", "1行=1チャプター (詳細メタ)"),
        ("chapter_topics.csv", "5,699", "Bridge", "1行=1チャプター×1トピック"),
    ]

    # ヘッダー行
    y = 1.8
    header_bg = TEAL_DARK
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(12), Inches(0.5),
              fill=header_bg)
    headers = [("ファイル名", 0.85, 3.0), ("行数", 3.95, 1.0), ("役割", 5.05, 2.2), ("説明", 7.35, 5.3)]
    for text, x, w in headers:
        add_text(slide, Inches(x), Inches(y), Inches(w), Inches(0.5),
                 text, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    y += 0.5
    for i, (fname, cnt, role, desc) in enumerate(files):
        # zebra
        row_bg = CREAM if i % 2 == 0 else WHITE
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(12), Inches(0.7),
                  fill=row_bg, line=BORDER)
        add_text(slide, Inches(0.85), Inches(y), Inches(3.0), Inches(0.7),
                 fname, size=12, bold=True, color=TEAL_DARK, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(3.95), Inches(y), Inches(1.0), Inches(0.7),
                 cnt, size=13, bold=True, color=ORANGE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, Inches(5.05), Inches(y), Inches(2.2), Inches(0.7),
                 role, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(7.35), Inches(y), Inches(5.3), Inches(0.7),
                 desc, size=11, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.7

    add_footer(slide)


def _draw_table_box(slide, x, y, w, h, title, columns, *, is_fact=False):
    """1つのテーブルボックスを描画。columnsは主要カラム名だけ数個。"""
    # ヘッダー
    header_color = ORANGE if is_fact else TEAL
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.4),
                       fill=header_color)
    add_text(slide, Inches(x), Inches(y), Inches(w), Inches(0.4),
             title, size=11, bold=True, color=WHITE, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 本体
    body_h = 0.32 * len(columns) + 0.1
    body = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.4), Inches(w), Inches(body_h),
                     fill=WHITE, line=header_color, line_width=1.5)
    # カラム行
    for i, (colname, key_label) in enumerate(columns):
        row_y = y + 0.4 + 0.05 + 0.32 * i
        col_color = ORANGE if key_label == "PK" else (TEAL_DARK if key_label == "FK" else DARK)
        add_text(slide, Inches(x + 0.1), Inches(row_y), Inches(w - 0.9), Inches(0.28),
                 colname, size=10, color=col_color, font=FONT_MONO,
                 bold=(key_label in ("PK", "FK")), anchor=MSO_ANCHOR.MIDDLE)
        if key_label:
            add_text(slide, Inches(x + w - 0.8), Inches(row_y), Inches(0.7), Inches(0.28),
                     key_label, size=9, bold=True, color=col_color, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.RIGHT)


def _draw_line(slide, x1, y1, x2, y2, color=TEAL, width=1.5):
    """2点間を結ぶ線 (単純な水平/垂直/斜め)"""
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)


def slide_schema_diagram(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_slide_header(slide, "スタースキーマ図", "episodes.csv を中心に4つのテーブルが接続")

    # 中央: episodes.csv (Fact)
    fact_x, fact_y, fact_w = 5.0, 3.2, 3.5
    _draw_table_box(slide, fact_x, fact_y, fact_w, 1.8, "episodes.csv (Fact)", [
        ("ep", "PK"),
        ("date", ""),
        ("duration_sec", ""),
        ("tweet_count", ""),
        ("... 他11カラム", ""),
    ], is_fact=True)

    # 左上: episode_speakers.csv
    sp_x, sp_y, sp_w = 0.9, 1.7, 3.5
    _draw_table_box(slide, sp_x, sp_y, sp_w, 1.5, "episode_speakers.csv", [
        ("ep", "FK"),
        ("date", "FK"),
        ("speaker_name", ""),
        ("speaker_handle", ""),
        ("role", ""),
    ])
    _draw_line(slide, sp_x + sp_w, sp_y + 0.8, fact_x, fact_y + 0.5)

    # 左下: episode_topics.csv
    tp_x, tp_y, tp_w = 0.9, 5.0, 3.5
    _draw_table_box(slide, tp_x, tp_y, tp_w, 1.5, "episode_topics.csv", [
        ("ep", "FK"),
        ("date", "FK"),
        ("topic", ""),
    ])
    _draw_line(slide, tp_x + tp_w, tp_y + 0.6, fact_x, fact_y + 1.0)

    # 右上: chapters.csv
    ch_x, ch_y, ch_w = 9.4, 1.7, 3.5
    _draw_table_box(slide, ch_x, ch_y, ch_w, 1.5, "chapters.csv", [
        ("ep", "FK"),
        ("date", "FK"),
        ("chapter_num", "PK"),
        ("title, summary", ""),
        ("start/end_sec", ""),
    ])
    _draw_line(slide, ch_x, ch_y + 0.8, fact_x + fact_w, fact_y + 0.5)

    # 右下: chapter_topics.csv (Bridge, chapters と接続)
    ctp_x, ctp_y, ctp_w = 9.4, 5.0, 3.5
    _draw_table_box(slide, ctp_x, ctp_y, ctp_w, 1.5, "chapter_topics.csv", [
        ("ep", "FK"),
        ("date", "FK"),
        ("chapter_num", "FK"),
        ("topic", ""),
    ])
    _draw_line(slide, ctp_x + ctp_w / 2, ctp_y, ch_x + ch_w / 2, ch_y + 1.5)

    # 凡例
    add_text(slide, Inches(0.7), Inches(6.75), Inches(4), Inches(0.3),
             "PK = 主キー / FK = 外部キー", size=10, color=GRAY)
    add_text(slide, Inches(9.0), Inches(6.75), Inches(4), Inches(0.3),
             "オレンジ枠 = Fact / ティール枠 = Dim・Bridge", size=10, color=GRAY, align=PP_ALIGN.RIGHT)

    add_footer(slide)


def slide_schema_details(prs, table_name, filename, subtitle, columns):
    """1テーブルの詳細カラム定義スライド。
    columns: list of (name, type, description)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_slide_header(slide, table_name, subtitle)

    add_text(slide, Inches(0.7), Inches(1.25), Inches(12), Inches(0.4),
             filename, size=13, bold=True, color=ORANGE, font=FONT_MONO)

    # ヘッダー
    y = 1.85
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(12), Inches(0.45),
              fill=TEAL_DARK)
    headers = [("カラム名", 0.85, 3.0), ("型", 3.95, 1.5), ("説明", 5.55, 7.0)]
    for text, x, w in headers:
        add_text(slide, Inches(x), Inches(y), Inches(w), Inches(0.45),
                 text, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.45

    row_h = 0.42
    for i, (name, type_, desc) in enumerate(columns):
        row_bg = CREAM if i % 2 == 0 else WHITE
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(12), Inches(row_h),
                  fill=row_bg, line=BORDER)
        # PK/FK は名前を強調
        name_color = ORANGE if "PK" in desc[:6] else (TEAL_DARK if "FK" in desc[:6] else DARK)
        add_text(slide, Inches(0.85), Inches(y), Inches(3.0), Inches(row_h),
                 name, size=11, bold=True, color=name_color, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(3.95), Inches(y), Inches(1.5), Inches(row_h),
                 type_, size=10, color=GRAY, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(5.55), Inches(y), Inches(7.0), Inches(row_h),
                 desc, size=10, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h
        if y > 6.9:
            break

    add_footer(slide)


def slide_tableau_howto(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_slide_header(slide, "Tableau 接続手順", "Join でなく Relationship を使う")

    steps = [
        ("1", "episodes.csv をドラッグ", "データソースパネルの中央に配置。これがメインテーブル"),
        ("2", "episode_speakers.csv を追加", "リレーションシップダイアログで ep + date の複合キー設定"),
        ("3", "episode_topics.csv を同じく接続", "ep + date で結合"),
        ("4", "chapters.csv を接続", "ep + date で結合"),
        ("5", "chapter_topics.csv を chapters.csv に接続", "ep + date + chapter_num の3キー結合"),
    ]

    y = 1.7
    for num, title, desc in steps:
        # 番号バッジ
        add_shape(slide, MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.6), Inches(0.6),
                  fill=ORANGE)
        add_text(slide, Inches(0.9), Inches(y), Inches(0.6), Inches(0.6),
                 num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # タイトル
        add_text(slide, Inches(1.7), Inches(y), Inches(11), Inches(0.35),
                 title, size=15, bold=True, color=TEAL_DARK)
        # 説明
        add_text(slide, Inches(1.7), Inches(y + 0.32), Inches(11), Inches(0.3),
                 desc, size=11, color=GRAY)
        y += 0.85

    # 警告ボックス
    warn_y = 6.15
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(warn_y), Inches(12), Inches(0.85),
              fill=RGBColor(0xFF, 0xF5, 0xE1), line=ORANGE)
    add_text(slide, Inches(0.9), Inches(warn_y + 0.1), Inches(11.6), Inches(0.3),
             "⚠  Join ではなく Relationship を使ってください", size=13, bold=True, color=ORANGE)
    add_text(slide, Inches(0.9), Inches(warn_y + 0.45), Inches(11.6), Inches(0.35),
             "Long format のテーブルを Join すると集計値が重複します。Relationship なら Tableau が集計時に自動で粒度調整します。",
             size=10, color=DARK)

    add_footer(slide)


def slide_analysis_ideas(prs, title, ideas):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_slide_header(slide, "分析アイデア", title)

    y = 1.7
    for label, desc in ideas:
        # 色バー
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(0.15), Inches(0.9),
                  fill=TEAL)
        add_text(slide, Inches(1.0), Inches(y), Inches(11.7), Inches(0.4),
                 label, size=14, bold=True, color=TEAL_DARK)
        add_text(slide, Inches(1.0), Inches(y + 0.42), Inches(11.7), Inches(0.5),
                 desc, size=11, color=DARK)
        y += 1.05

    add_footer(slide)


def slide_caveats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_slide_header(slide, "注意事項・限界", "分析結果の解釈時に留意すべきこと")

    items = [
        ("文字起こし精度", "AI (Whisper) 自動処理のため、固有名詞・専門用語で誤認識あり。トピック抽出も LLM の解釈が入る"),
        ("発話者識別なし", "チャプター/トピックが「誰の発言か」は現状特定不可。speaker × topic は「その回に登場した話題」で、その人が話したとは限らない"),
        ("tweet 集計の抜け", "#Vizトーク ハッシュタグ付きのみ。実況していない回や、ハッシュタグ付け忘れのツイートは含まれない"),
        ("第35回以前", "音源が残っていない回はメタデータのみ (recording=0)。分析時は recording=1 でフィルタが基本"),
        ("トピックの正規化", "LLM が抽出したままなので、同義語 (例: Tableau Prep / Prep) の統合は未実施"),
    ]

    y = 1.7
    for label, desc in items:
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(0.15), Inches(0.85),
                  fill=ORANGE)
        add_text(slide, Inches(1.0), Inches(y), Inches(11.7), Inches(0.4),
                 label, size=14, bold=True, color=ORANGE)
        add_text(slide, Inches(1.0), Inches(y + 0.42), Inches(11.7), Inches(0.4),
                 desc, size=11, color=DARK)
        y += 1.0

    add_footer(slide)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    add_text(slide, Inches(0), Inches(2.3), Inches(13.33), Inches(0.8),
             "データセットは以上です", size=32, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.665), Inches(3.2), Inches(2.0), Inches(0.03),
              fill=ORANGE)

    add_text(slide, Inches(0), Inches(3.5), Inches(13.33), Inches(0.5),
             "質問・要望・共同分析の相談は主催者陣まで", size=15, color=DARK, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0), Inches(4.5), Inches(13.33), Inches(0.5),
             "データ更新: python3 bin/export_analytics.py",
             size=12, color=GRAY, align=PP_ALIGN.CENTER, font=FONT_MONO)

    add_text(slide, Inches(0), Inches(6.5), Inches(13.33), Inches(0.4),
             "Enjoy your analysis!  📊", size=18, color=ORANGE, align=PP_ALIGN.CENTER)


# =========================================================
# メイン
# =========================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)   # 16:9
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_overview(prs)
    slide_files_list(prs)
    slide_schema_diagram(prs)

    slide_schema_details(prs, "episodes.csv — メイン Fact テーブル", "episodes.csv",
        "1行 = 1回。分析の起点になるテーブル",
        [
            ("ep", "int", "PK 回番号 (36, 37, ...)"),
            ("date", "date", "開催日 YYYY-MM-DD"),
            ("year", "int", "開催年"),
            ("month", "int", "開催月 (1-12)"),
            ("weekday", "string", "曜日 (月/火/水/木/金/土/日)"),
            ("title", "string", "オリジナルのスペースタイトル"),
            ("recording", "int (0/1)", "録音が残っているか"),
            ("duration_sec / min / hour", "int/float", "放送時間 (3つの単位で提供)"),
            ("speaker_count", "int", "出演スピーカー数"),
            ("topic_count", "int", "この回の主要トピック数"),
            ("chapter_count", "int", "この回のチャプター数"),
            ("tweet_count", "int", "#Vizトーク の実況ツイート数"),
        ])

    slide_schema_details(prs, "episode_speakers.csv — Speaker Dim", "episode_speakers.csv",
        "1行 = 1回×1スピーカー (long format)",
        [
            ("ep", "int", "FK → episodes.ep"),
            ("date", "date", "FK 補助キー"),
            ("speaker_name", "string", "スピーカー表示名"),
            ("speaker_handle", "string", "X ハンドル (@なし)"),
            ("role", "string", "ホスト / 共同ホスト / スピーカー"),
        ])

    slide_schema_details(prs, "episode_topics.csv — Topic Dim", "episode_topics.csv",
        "1行 = 1回×1トピック (long format)",
        [
            ("ep", "int", "FK → episodes.ep"),
            ("date", "date", "FK 補助キー"),
            ("topic", "string", "トピック名 (例: LOD計算, Tableau Prep, 雑談)"),
        ])

    slide_schema_details(prs, "chapters.csv — Chapter Dim", "chapters.csv",
        "1行 = 1チャプター (話題の切り替わり単位)",
        [
            ("ep", "int", "FK → episodes.ep"),
            ("date", "date", "FK 補助キー"),
            ("chapter_num", "int", "その回のチャプター通し番号 (1〜)"),
            ("start_time", "H:MM:SS", "開始時刻 (音源内)"),
            ("start_sec", "int", "開始秒数"),
            ("end_time / end_sec", "H:MM:SS / int", "終了時刻・秒"),
            ("duration_sec", "int", "チャプターの長さ (秒)"),
            ("title", "string", "LLM 生成のチャプター見出し"),
            ("summary", "string", "LLM 生成のチャプター要約"),
            ("tag_count", "int", "このチャプターに付与されたトピック数"),
        ])

    slide_schema_details(prs, "chapter_topics.csv — Chapter × Topic Bridge", "chapter_topics.csv",
        "1行 = 1チャプター×1トピック (chapters と結合して使う)",
        [
            ("ep", "int", "FK → chapters.ep"),
            ("date", "date", "FK 補助キー"),
            ("chapter_num", "int", "FK → chapters.chapter_num"),
            ("start_sec", "int", "チャプター開始秒 (chaptersから引き継ぎ)"),
            ("topic", "string", "トピック名"),
        ])

    slide_tableau_howto(prs)

    slide_analysis_ideas(prs, "スピーカー観点", [
        ("出演回数ランキング", "episode_speakers + episodes で COUNTD(ep)。ホスト陣とスピーカーで分離できる"),
        ("スピーカー別の平均トーク時間・平均tweet数", "episodes.duration_min / tweet_count と結合"),
        ("初出演 → 最終出演の期間", "MIN(date) と MAX(date) で活動期間を測る"),
        ("よく一緒に出るスピーカーのネットワーク", "episode_speakers の自己結合、共演回数マトリクス"),
    ])

    slide_analysis_ideas(prs, "トピック観点", [
        ("トピック別の出現頻度・時系列トレンド", "episode_topics × date の年月ヒートマップ"),
        ("話題のライフサイクル", "あるトピックの初登場〜衰退までのタイムライン"),
        ("Tableau 新機能リリースと話題の相関", "外部データ (製品リリース日) と重ねる"),
        ("チャプター単位の話題遷移", "chapter_topics で1回内の話題の変化を見る"),
    ])

    slide_analysis_ideas(prs, "番組全体・クロス集計", [
        ("開催頻度・時間の推移", "date × duration_min の折れ線 (年月別 or 週別)"),
        ("ツイート数と話題の相関", "tweet_count と topic を組み合わせ、バズった話題を特定"),
        ("チャプター数と放送時間の関係", "話題転換の速さを測る (chapter_count / duration_hour)"),
        ("スピーカー × トピックのヒートマップ", "「誰が何を語るか」の可視化。ただし発話者未識別の限界あり"),
    ])

    slide_caveats(prs)
    slide_closing(prs)

    OUT_PATH.parent.mkdir(exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"✅ 生成完了: {OUT_PATH}")
    print(f"   スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
